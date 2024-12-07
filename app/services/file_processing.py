from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

class FileProcessing:
    def __init__(self):
        pass

    def load_word(file_path):
        doc = DocxDocument(file_path)
        text = [para.text for para in doc.paragraphs]
        return "\n".join(text)

    def load_excel(file_path):
        wb = load_workbook(file_path)
        text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                row_text = " ".join(str(cell) for cell in row if cell)
                text.append(row_text)
        return "\n".join(text)

    def load_ppt(file_path):
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text.append(paragraph.text)
        return "\n".join(text)

    def load_office_file(file_path):
        if file_path.endswith(".docx"):
            content = load_word(file_path)
        elif file_path.endswith(".xlsx"):
            content = load_excel(file_path)
        elif file_path.endswith(".pptx"):
            content = load_ppt(file_path)
        else:
            raise ValueError("Unsupported file type")
        return [Document(page_content=content)]
